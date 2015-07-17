#!/usr/bin/env python
# encoding: utf-8
"""
Generate Markov chain given set of files.
"""
import codecs
import locale
import logging
import optparse
import os
import sys
import warnings
import os.path
import docx
import markovify
import pyPdf
import html2text
import pptx

__author__ = 'Simon Brunning'
__version__ = "0.1"

script_name = os.path.basename(sys.argv[0])
usage = script_name + ' [options] source-files'
description = '''
Generate Markov chain given set of files.
'''

logger = logging.getLogger(script_name)


def main(*argv):
    options, script, args, help = get_options(argv)
    init_logger(options.verbosity)

    corpus = ''
    for arg in args:
        ext = os.path.splitext(arg)[-1]
        strategy = extract_strategies.get(ext, UnknownFiletypeStrategy)
        corpus += strategy(arg).extract()

    text_model = markovify.Text(corpus)

    # Print five randomly-generated sentences
    for i in range(options.length):
        print(text_model.make_sentence())


class ExtractStrategy(object):
    def __init__(self, filename):
        self.filename = filename
        logger.info("Adding contents of %s", filename)

    def extract(self):
        raise NotImplementedError()


class DocxStrategy(ExtractStrategy):
    def extract(self):
        document = docx.Document(self.filename)
        return " ".join(paragraph.text for paragraph in document.paragraphs)


class PptxStrategy(ExtractStrategy):
    def extract(self):
        """From http://python-pptx.readthedocs.org/en/latest/user/quickstart.html#extract-all-text-from-slides-in-presentation"""

        prs = pptx.Presentation(self.filename)

        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return " ".join(text_runs)


class PdfStrategy(ExtractStrategy):
    def extract(self):
        pdf = pyPdf.PdfFileReader(open(self.filename, "rb"))
        return " ".join(page.extractText() for page in pdf.pages)


class HtmlStrategy(ExtractStrategy):
    def extract(self):
        with open(self.filename) as html:
            return html2text.html2text(html.read())


class TxtStrategy(ExtractStrategy):
    def extract(self):
        with open(self.filename) as f:
            return f.read()


class UnknownFiletypeStrategy(ExtractStrategy):
    def __init__(self, filename):
        logger.warn("Unknown file type %s", filename)

    def extract(self):
        return ""


extract_strategies = {
    ".docx": DocxStrategy,
    ".pptx": PptxStrategy,
    ".pdf": PdfStrategy,
    ".html": HtmlStrategy,
    ".txt": TxtStrategy,
    ".md": TxtStrategy,
}


def get_options(argv):
    '''Get options and arguments from argv string.'''
    parser = optparse.OptionParser(usage=usage, version=__version__)
    parser.description = description
    parser.add_option("-v", "--verbosity", action="count", default=0,
                      help="Specify up to three times to increase verbosity, i.e. -v to see warnings, "
                           "-vv for information messages, or -vvv for debug messages.")
    parser.add_option("-l", "--length", type="int", help="Length (in sentences.)", default=5)

    options, args = parser.parse_args(list(argv))
    script, args = args[0], args[1:]
    return options, script, args, parser.format_help()


def init_logger(verbosity, stream=sys.stdout):
    '''Initialize logger and warnings according to verbosity argument.
    Verbosity levels of 0-3 supported.'''
    is_not_debug = verbosity <= 2
    level = [logging.ERROR, logging.WARNING, logging.INFO][verbosity] if is_not_debug else logging.DEBUG
    format = '%(message)s' if is_not_debug \
        else '%(asctime)s %(levelname)-8s %(name)s %(module)s.py:%(funcName)s():%(lineno)d %(message)s'
    logging.basicConfig(level=level, format=format, stream=stream)
    if is_not_debug: warnings.filterwarnings('ignore')


def wrap_stream_for_tty(stream):
    if stream.isatty():
        # Configure locale from the user's environment settings.
        locale.setlocale(locale.LC_ALL, '')

        # Wrap stdout with an encoding-aware writer.
        lang, encoding = locale.getdefaultlocale()
        logger.debug('Streaming to tty with lang, encoding = %s, %s', lang, encoding)
        if encoding:
            return codecs.getwriter(encoding)(stream)
        else:
            logger.warn('No tty encoding found!')

    return stream


if __name__ == "__main__":
    sys.exit(main(*sys.argv))
